import os
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import fitz  # PyMuPDF
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_huggingface import HuggingFaceEmbeddings

# Resolve data folder path relative to this script    { CHANGED THE PATH TO DATA FOLDER WHICH MAKES IT EASY FOR OTHERS.}
DATA_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__),"data"))

if not os.path.exists(DATA_DIR):
    raise FileNotFoundError(f"DATA_DIR not found at: {DATA_DIR}")

def extract_text_from_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        if "question" in df.columns and "answer" in df.columns:
            return (df["question"].astype(str) + " " + df["answer"].astype(str)).tolist()
        else:
            print(f"Skipping {file_path}: missing 'question' or 'answer' column")
            return []
    except Exception as e:
        print(f"Error reading CSV: {e}")
        return []

def extract_text_from_ppt_with_ocr(file_path):
    prs = Presentation(file_path)
    extracted_texts = []

    for slide in prs.slides:
        slide_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + "\n"

        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        slide_text += "\n[Image OCR Text]:\n" + ocr_text + "\n"
                except Exception as e:
                    print(f"OCR failed on image in PPT: {e}")

        if slide_text.strip():
            extracted_texts.append(slide_text.strip())

    return extracted_texts

def extract_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return [f.read().strip()]
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return []

def extract_text_from_pdf(file_path):
    extracted_texts = []
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            images = page.get_images(full=True)
            page_combined_text = page_text.strip()

            for img_index, img in enumerate(images):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    img_pil = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(img_pil)
                    if ocr_text.strip():
                        page_combined_text += f"\n[Image OCR Page {page_num+1} Img {img_index+1}]:\n" + ocr_text
                except Exception as e:
                    print(f"OCR failed on image in PDF page {page_num+1}: {e}")

            if page_combined_text.strip():
                extracted_texts.append(page_combined_text)

    except Exception as e:
        print(f"Error reading PDF: {e}")
    
    return extracted_texts

# Gather all text blocks
all_text_blocks = []

for filename in os.listdir(DATA_DIR):
    file_path = os.path.join(DATA_DIR, filename)

    if filename.endswith(".csv"):
        print(f"Reading CSV: {filename}")
        all_text_blocks.extend(extract_text_from_csv(file_path))
    elif filename.endswith(".pptx"):
        print(f"Reading PPTX (with OCR): {filename}")
        all_text_blocks.extend(extract_text_from_ppt_with_ocr(file_path))
    elif filename.endswith(".txt"):
        print(f"Reading TXT: {filename}")
        all_text_blocks.extend(extract_text_from_txt(file_path))
    elif filename.endswith(".pdf"):
        print(f"Reading PDF (with OCR): {filename}")
        all_text_blocks.extend(extract_text_from_pdf(file_path))
    else:
        print(f"⏭ Skipping unsupported file: {filename}")

# Embed and save to Chroma
print("Creating vector DB...")
embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
vectordb = Chroma.from_texts(texts=all_text_blocks, embedding=embedding_model, persist_directory="chroma_db")
vectordb.persist()

print("Vector DB updated with all files in /data.")
print("Done! You can now query the vector database.")

#new line  
print(f"Indexed {len(all_text_blocks)} documents/text blocks.")

# Clean up


